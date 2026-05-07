"""
Apply Winter brand styling to slide 2 of `eICU based 150 P validation.pptx`.

Slide 1 is the Winter master template; slide 2 is our validation page.
We extract Winter's typography + palette from slide 1 and remap slide 2 to
match - without restructuring its layout or content geometry.

Winter brand DNA (lifted from slide 1):
    Fonts:   Segoe UI (regular)  /  Segoe UI Semibold (emphasis)
    Colors:  #454545 - dark    (big titles, metric numbers, bullet heads)
             #506370 - slate   (section heads, semibold accents)
             #898A87 - light   (body, captions, decorative chars)
    Soft:    #F2F2F0 - cream   (Winter ribbon fill / pill fill)

Idempotent: safe to re-run.

    python presentations/scripts/apply_winter_brand_to_validation.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[2]
DECK = ROOT / "presentations" / "current" / "eICU based 150 P validation.pptx"

# Winter palette
DARK = RGBColor(0x45, 0x45, 0x45)
SLATE = RGBColor(0x50, 0x63, 0x70)
LIGHT = RGBColor(0x89, 0x8A, 0x87)
CREAM = RGBColor(0xF2, 0xF2, 0xF0)

# Decorative single-char runs (arrows / bullets) - always rendered LIGHT.
DECORATIVE_CHARS = set("\u2192\u2190\u2193\u2191\u2022\u00b7\u25e6\u25b8\u25b9\u2014\u2013\u2192")


def compute_winter_color(text: str, size: float,
                         font_name: str, is_bold: bool) -> RGBColor:
    """Map a text run to the Winter palette based on its inferred role.

    Visual hierarchy from slide 1 (Winter master):
        size >= 20                       -> DARK   (big titles, big numbers)
        size >= 11 and emphasis          -> SLATE  (section / column heads)
        emphasis (smaller)               -> DARK   (bullet heads / pill titles)
        regular                          -> LIGHT  (body, captions)
        decorative single-char           -> LIGHT  (arrows / bullets)

    "emphasis" = bold attribute OR font name contains "Semibold". The font-name
    fallback makes this idempotent across re-runs (after we move emphasis from
    the bold attribute to the Semibold font).
    """
    if text and all(c in DECORATIVE_CHARS for c in text):
        return LIGHT

    is_emphasis = is_bold or "Semibold" in (font_name or "")

    if size >= 20:
        return DARK
    if size >= 11 and is_emphasis:
        return SLATE
    if is_emphasis:
        return DARK
    return LIGHT


def restyle_run(run):
    """Snapshot original state -> compute color -> mutate font + apply color.
    The order matters: we must read bold/font BEFORE clearing bold attribute."""
    text = (run.text or "").strip()
    orig_bold = bool(run.font.bold)
    orig_name = run.font.name
    size = run.font.size.pt if run.font.size else 11

    color = compute_winter_color(text, size, orig_name, orig_bold)
    is_emphasis = orig_bold or "Semibold" in (orig_name or "")

    # Big titles / metric numbers: regular Segoe UI - matches Winter's
    # "Unstructured" / "Deterministic" / "Per-Hospital" treatment (24pt regular).
    if size >= 20:
        run.font.name = "Segoe UI"
        run.font.bold = False
    elif is_emphasis:
        # Emphasis -> Semibold via font name, NOT bold attribute (Winter style).
        run.font.name = "Segoe UI Semibold"
        run.font.bold = False
    else:
        run.font.name = "Segoe UI"
        run.font.bold = False

    run.font.color.rgb = color


def restyle_text(slide) -> int:
    n = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text.strip():
                    continue
                restyle_run(run)
                n += 1
    return n


def restyle_pills(slide) -> int:
    """The 5 pipeline pills currently use vivid coloured fills/borders.
    Winter uses one soft cream fill across all panels."""
    n = 0
    for shape in slide.shapes:
        if "[guardrail-features-pill]" not in shape.name:
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            shape.fill.solid()
            shape.fill.fore_color.rgb = CREAM
            shape.line.fill.background()
            n += 1
    return n


def main():
    if not DECK.exists():
        raise SystemExit(f"Deck not found: {DECK}")

    prs = Presentation(str(DECK))
    if len(prs.slides) < 2:
        raise SystemExit("Deck has fewer than 2 slides")

    slide = prs.slides[1]

    n_runs = restyle_text(slide)
    n_pills = restyle_pills(slide)

    print(f"Restyled {n_runs} text runs to Winter typography + palette")
    print(f"Restyled {n_pills} pipeline pills to Winter cream fill")

    try:
        prs.save(str(DECK))
        print(f"\nSaved: {DECK}")
    except PermissionError:
        alt = DECK.with_name(DECK.stem + "_winter.pptx")
        prs.save(str(alt))
        print(f"\nDeck is locked (open in PowerPoint).")
        print(f"Saved alternate copy to: {alt}")
        print("Please close PowerPoint, then either replace the original "
              "with the alternate, or re-run this script.")


if __name__ == "__main__":
    main()
