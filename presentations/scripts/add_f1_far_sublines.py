"""
Add the missing sub-lines under the F1 score and False-alarm rate metrics
on slide 17 of Sepsis_GenAI_DeepDivev7.pptx, matching the existing pattern
used by Sensitivity and Specificity (count + 95% CI).

Computed from the actual v7 confusion matrix (TP=28, FP=43, TN=73, FN=6):
    F1            = 0.533       Bootstrap 95% CI = 0.42 - 0.64
                                Components: Precision 39.4% x Sens 82.4%
    False-alarm   = 37.1%       Wilson 95% CI    = 28.8 - 46.1
                                Components: 43 of 116 controls

Idempotent: re-running first deletes any sub-lines this script previously added.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
DECK = ROOT / "presentations" / "current" / "Sepsis_GenAI_DeepDivev7.pptx"

# Match the existing v7 deck palette (Calibri / blue-slate) - NOT Winter.
DARK = RGBColor(0x1F, 0x2F, 0x3F)
SLATE = RGBColor(0x55, 0x66, 0x77)

OWN_TAG = "[f1-far-subline]"


def add_text(slide, left_in, top_in, width_in, height_in, text,
             *, size=10, bold=False, color=DARK, font="Calibri"):
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in),
                                   Inches(width_in), Inches(height_in))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    box.name = f"{OWN_TAG}:{box.name}"
    return box


def remove_owned(slide) -> int:
    spTree = slide.shapes._spTree
    n = 0
    for sh in list(slide.shapes):
        if sh.name and OWN_TAG in sh.name:
            spTree.remove(sh._element)
            n += 1
    return n


def main():
    if not DECK.exists():
        raise SystemExit(f"Deck not found: {DECK}")

    prs = Presentation(str(DECK))
    slide = prs.slides[-1]  # validation milestone slide

    n_removed = remove_owned(slide)
    if n_removed:
        print(f"Removed {n_removed} previously-added sub-lines (idempotent).")

    # Existing sub-line geometry under Sensitivity/Specificity:
    #   line 1 (count):  T=2.72, H=0.25, size 9, color #1F2F3F (DARK)
    #   line 2 (95% CI): T=2.98, H=0.25, size 8, color #556677 (SLATE)

    # F1 score column - L=4.85
    add_text(slide, 4.85, 2.72, 1.95, 0.25,
             "PPV 39.4%  \u00d7  Sens 82.4%",
             size=9, color=DARK)
    add_text(slide, 4.85, 2.98, 1.95, 0.25,
             "95% CI 0.42 \u2013 0.64",
             size=8, color=SLATE)

    # False-alarm rate column - L=6.85
    add_text(slide, 6.85, 2.72, 1.95, 0.25,
             "43 of 116 controls",
             size=9, color=DARK)
    add_text(slide, 6.85, 2.98, 1.95, 0.25,
             "95% CI 28.8 \u2013 46.1",
             size=8, color=SLATE)

    print("Added 2 sub-lines under F1 score:")
    print("    'PPV 39.4% \u00d7 Sens 82.4%'")
    print("    '95% CI 0.42 \u2013 0.64'")
    print("Added 2 sub-lines under False-alarm rate:")
    print("    '43 of 116 controls'")
    print("    '95% CI 28.8 \u2013 46.1'")

    try:
        prs.save(str(DECK))
        print(f"\nSaved: {DECK}")
    except PermissionError:
        alt = DECK.with_name(DECK.stem + "_subs.pptx")
        prs.save(str(alt))
        print(f"\nDeck is open in PowerPoint - saved alternate to {alt}")


if __name__ == "__main__":
    main()
