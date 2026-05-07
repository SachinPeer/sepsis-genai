"""
Append a one-page "Validation Milestone" slide to Sepsis_GenAI_DeepDivev6.pptx
and save as v7. Matches the visual language of slides 15-16 (Calibri body /
24-pt blue title / 22-pt big metric numbers / coloured section headers).

Run:
    python presentations/scripts/append_validation_milestone_slide.py
"""
from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# -- Paths ----------------------------------------------------------------
ROOT = Path(__file__).resolve().parents[2]
SOURCE = ROOT / "presentations" / "current" / "Sepsis_GenAI_DeepDivev6.pptx"
TARGET = ROOT / "presentations" / "current" / "Sepsis_GenAI_DeepDivev7.pptx"

# -- Brand palette (matches slides 15-16) ---------------------------------
BLUE = RGBColor(0x1E, 0x6C, 0x93)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
ORANGE = RGBColor(0xE6, 0x5C, 0x00)
SLATE = RGBColor(0x55, 0x66, 0x77)
DARK = RGBColor(0x1F, 0x2F, 0x3F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
LIGHT_BLUE = RGBColor(0xE3, 0xF2, 0xFD)
LIGHT_ORANGE = RGBColor(0xFF, 0xF3, 0xE0)
LIGHT_RED = RGBColor(0xFF, 0xEB, 0xEE)
RAIL_BG = RGBColor(0x1E, 0x6C, 0x93)


# -- Helpers --------------------------------------------------------------
def add_text(slide, left_in, top_in, width_in, height_in, text,
             *, size=10, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             font="Calibri"):
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in),
                                   Inches(width_in), Inches(height_in))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_filled_box(slide, left_in, top_in, width_in, height_in,
                   fill, *, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(left_in), Inches(top_in),
                                Inches(width_in), Inches(height_in))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.5)
    sh.shadow.inherit = False
    return sh


def add_rounded_box(slide, left_in, top_in, width_in, height_in,
                    fill, *, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(left_in), Inches(top_in),
                                Inches(width_in), Inches(height_in))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.5)
    sh.shadow.inherit = False
    return sh


def metric_tile(slide, left, top, width, label, value, sub, ci,
                value_color=BLUE):
    add_text(slide, left, top, width, 0.28, label,
             size=10, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + 0.27, width, 0.55, value,
             size=22, bold=True, color=value_color, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + 0.87, width, 0.25, sub,
             size=9, color=DARK, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + 1.13, width, 0.25, ci,
             size=8, color=SLATE, align=PP_ALIGN.CENTER)


def bullet(slide, left, top, width, lead, body, lead_color=DARK):
    """Bold lead phrase with a sub-line of supporting detail."""
    add_text(slide, left, top, width, 0.22, "•  " + lead,
             size=10, bold=True, color=lead_color)
    add_text(slide, left + 0.20, top + 0.22, width - 0.20, 0.32, body,
             size=9, color=SLATE)


# -- Slide builder --------------------------------------------------------
def build_milestone_slide(prs):
    blank_layout = prs.slide_layouts[6]  # "Blank"
    slide = prs.slides.add_slide(blank_layout)

    # ===== TITLE BAR =====
    add_text(slide, 0.69, 0.27, 12.0, 0.50,
             "Validation Milestone — eICU v4 | C1 + C2 Reasoning-Aware Guardrail",
             size=24, bold=True, color=BLUE)
    add_text(slide, 0.80, 0.85, 12.0, 0.32,
             "150 PHI-compliant eICU-CRD patients (34 sepsis / 116 controls) "
             "• Real ICU nurse notes • Claude Sonnet 4.5 @ temperature=0 (deterministic) "
             "• Same patient \u2192 same answer, every time",
             size=10, color=SLATE)

    # ===== ROW 1: HEADLINE METRICS (4 big tiles) =====
    metrics_top = 1.40
    add_text(slide, 0.85, metrics_top, 12.0, 0.30,
             "Headline result — locked, byte-reproducible",
             size=13, bold=True, color=BLUE)

    metric_tile(slide, 0.85, metrics_top + 0.45, 1.85,
                "Sensitivity", "82.4%", "28 of 34 sepsis",
                "95% CI 66.5 – 91.7", value_color=GREEN)
    metric_tile(slide, 2.85, metrics_top + 0.45, 1.85,
                "Specificity", "62.9%", "73 of 116 controls",
                "95% CI 53.9 – 71.2", value_color=BLUE)
    metric_tile(slide, 4.85, metrics_top + 0.45, 1.85,
                "F1 score", "0.533", "from 0.392 baseline",
                "+36% improvement", value_color=GREEN)
    metric_tile(slide, 6.85, metrics_top + 0.45, 1.85,
                "False-alarm rate", "37.1%", "from 69.8% baseline",
                "−32.7 pp", value_color=GREEN)

    # ===== TOP-RIGHT: Salient sepsis catches =====
    sal_left = 8.95
    sal_top = metrics_top + 0.45
    add_rounded_box(slide, sal_left, sal_top, 3.85, 1.40,
                    LIGHT_GREEN, line=GREEN)
    add_text(slide, sal_left + 0.12, sal_top + 0.06, 3.65, 0.28,
             "7 sepsis cases caught that bedside SIRS / qSOFA "
             "would have MISSED",
             size=10, bold=True, color=GREEN)
    add_text(slide, sal_left + 0.12, sal_top + 0.34, 3.65, 0.22,
             "(qSOFA < 2 AND SIRS < 2 — formal screening would say \"no sepsis\")",
             size=8, color=SLATE)
    add_text(slide, sal_left + 0.12, sal_top + 0.58, 3.65, 0.22,
             "•  p00027 \u2013 \"silent sepsis\": lactate + neurological decline",
             size=9, bold=True, color=DARK)
    add_text(slide, sal_left + 0.12, sal_top + 0.78, 3.65, 0.22,
             "•  p00008 \u2013 cryptic shock pattern (lactate + hypoperfusion)",
             size=9, bold=True, color=DARK)
    add_text(slide, sal_left + 0.12, sal_top + 0.98, 3.65, 0.22,
             "•  p00009 \u2013 compensated septic shock on norepinephrine",
             size=9, bold=True, color=DARK)
    add_text(slide, sal_left + 0.12, sal_top + 1.18, 3.65, 0.22,
             "•  p00021 \u2013 already in septic shock with end-organ dysfunction",
             size=9, bold=True, color=DARK)

    # ===== ROW 2: TWO-COLUMN BODY (architecture pillars + suppressor logic) =====
    row2_top = 3.50

    # ---- LEFT COLUMN: WHAT MAKES THIS UNIQUE ----
    add_text(slide, 0.80, row2_top, 5.95, 0.28,
             "What makes the model unique — \"reasoning-aware guardrail\"",
             size=12, bold=True, color=BLUE)

    bullet(slide, 0.85, row2_top + 0.34, 5.95,
           "Nurse-note narrative \u2192 LLM",
           "Unstructured shift notes ingested as primary signal "
           "(not just numbers). Captures the bedside story.")
    bullet(slide, 0.85, row2_top + 0.92, 5.95,
           "Deterministic clinical scores (qSOFA, SIRS, SOFA)",
           "Computed in code from vitals + labs and FED INTO the prompt — "
           "scores are calculated, not hallucinated.")
    bullet(slide, 0.85, row2_top + 1.50, 5.95,
           "C1 LLM-aware suppression",
           "When the LLM explicitly identifies a NON-sepsis cause "
           "(post-op stress, ACS, DKA, GI bleed, residual sedation, alkalosis), "
           "the guardrail RESPECTS that verdict instead of over-ruling it.")
    bullet(slide, 0.85, row2_top + 2.18, 5.95,
           "C2 pattern suppressor — \"don't trigger if LLM has explained it\"",
           "7-branch deterministic layer that holds back alerts when the LLM "
           "rationale fits a non-infectious archetype AND no rescue signal "
           "(lactate, GCS, septic-shock language, qSOFA \u2265 2, SOFA \u2265 4) is present.")

    # ---- RIGHT COLUMN: COMPLIANCE + AUDIT + DETERMINISM ----
    add_text(slide, 7.10, row2_top, 5.85, 0.28,
             "Production-ready pillars",
             size=12, bold=True, color=BLUE)

    bullet(slide, 7.15, row2_top + 0.34, 5.85,
           "Temperature = 0 (fully deterministic stack)",
           "Identical patient \u2192 identical risk + reasoning. "
           "Simulator and production agree to the patient.")
    bullet(slide, 7.15, row2_top + 0.92, 5.85,
           "PHI / HIPAA compliant by design",
           "eICU validation on open de-identified data. Production: AWS BAA, "
           "MongoDB BAA, encryption in transit + at rest, configurable retention.")
    bullet(slide, 7.15, row2_top + 1.50, 5.85,
           "Full audit log per prediction",
           "LLM's initial risk + priority, every guardrail decision, "
           "C1/C2 branch + reason, override triggers — all persisted.")
    bullet(slide, 7.15, row2_top + 2.18, 5.85,
           "Hospital-configurable guardrail (JSON, no code change)",
           "Each site can tune thresholds, denial phrases, "
           "rescue signals via genai_clinical_guardrail.json + UI editor.")

    # ===== BOTTOM TRAJECTORY STRIP =====
    strip_top = 6.65
    add_filled_box(slide, 0.50, strip_top, 12.30, 0.34, RAIL_BG)
    add_text(slide, 0.70, strip_top + 0.05, 12.0, 0.24,
             "Trajectory:  v4 baseline 30.2%  \u2192  v6 (C1 + scores in prompt) 39.7%  "
             "\u2192  v7 (C1 + scores + C2 suppressor) 62.9%  specificity     |     "
             "Sensitivity locked at 82.4% throughout — ZERO true-positive loss",
             size=10, bold=True, color=WHITE)

    # ===== FOOTER =====
    add_text(slide, 0.50, 7.20, 6.0, 0.25,
             "Sepsis GenAI  |  Medbeacon",
             size=9, color=SLATE)
    add_text(slide, 9.00, 7.20, 4.0, 0.25,
             "eICU-CRD v2.0.1 demo  |  validated 2026-02-11",
             size=9, color=SLATE, align=PP_ALIGN.RIGHT)


def main():
    if not SOURCE.exists():
        raise SystemExit(f"Source deck not found: {SOURCE}")

    if TARGET.exists():
        TARGET.unlink()
    shutil.copyfile(SOURCE, TARGET)

    prs = Presentation(str(TARGET))
    n_before = len(prs.slides)
    build_milestone_slide(prs)
    n_after = len(prs.slides)
    prs.save(str(TARGET))

    print(f"Source : {SOURCE}")
    print(f"Target : {TARGET}")
    print(f"Slides : {n_before} \u2192 {n_after} (added 1)")
    print(f"New slide title: \"Validation Milestone — eICU v4 | C1 + C2 Reasoning-Aware Guardrail\"")


if __name__ == "__main__":
    main()
