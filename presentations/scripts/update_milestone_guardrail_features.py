"""
Update the v7 milestone slide IN-PLACE.

Replaces the bottom area with a clean, non-technical 5-stage pipeline:

    Data Preprocess -> Scores -> LLM -> LLM-aware Guardrails -> Noisy-alert
    Suppression

Idempotent: each run removes anything we previously added (any 3-pill rows,
old "Trajectory" / "Guardrail elevates" strip) and re-builds the pipeline.

    python presentations/scripts/update_milestone_guardrail_features.py
"""
from __future__ import annotations

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
DECK = ROOT / "presentations" / "current" / "Sepsis_GenAI_DeepDivev7.pptx"

# Brand palette
BLUE = RGBColor(0x1E, 0x6C, 0x93)
DARK = RGBColor(0x1F, 0x2F, 0x3F)
SLATE = RGBColor(0x55, 0x66, 0x77)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# 5 stage colours - light fills with matching darker title/border
STAGE_FILL = [
    RGBColor(0xEC, 0xEF, 0xF1),  # Data Preprocess - cool gray
    RGBColor(0xE3, 0xF2, 0xFD),  # Scores          - light blue
    RGBColor(0xEE, 0xE6, 0xF5),  # LLM             - light purple
    RGBColor(0xFF, 0xF3, 0xE0),  # Guardrails      - light orange
    RGBColor(0xE8, 0xF5, 0xE9),  # Suppression     - light green
]
STAGE_BORDER = [
    RGBColor(0x60, 0x77, 0x88),
    RGBColor(0x1E, 0x6C, 0x93),
    RGBColor(0x6A, 0x1B, 0x9A),
    RGBColor(0xE6, 0x5C, 0x00),
    RGBColor(0x2E, 0x7D, 0x32),
]

# Tag we put on shapes we own, so re-runs can clean up cleanly.
OWN_TAG = "[guardrail-features-pill]"


# ----------------------------------------------------------------- helpers
def add_text(slide, left_in, top_in, width_in, height_in, text,
             *, size=10, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="Calibri",
             tag: Optional[str] = OWN_TAG):
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in),
                                   Inches(width_in), Inches(height_in))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if tag:
        box.name = f"{tag}:{box.name}"
    return box


def add_rounded_box(slide, left_in, top_in, width_in, height_in,
                    fill, *, line=None, tag: Optional[str] = OWN_TAG):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(left_in), Inches(top_in),
                                Inches(width_in), Inches(height_in))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    if tag:
        sh.name = f"{tag}:{sh.name}"
    return sh


# ------------------------------------------------------------- removal ops
def remove_owned_shapes(slide) -> int:
    spTree = slide.shapes._spTree
    removed = 0
    for sh in list(slide.shapes):
        if sh.name and OWN_TAG in sh.name:
            spTree.remove(sh._element)
            removed += 1
    return removed


def remove_old_strip(slide) -> int:
    """Remove the old blue trajectory / 'Guardrail elevates' strip if present.

    Two shapes make up the original strip:
      * a wide filled rectangle at T~6.65, L~0.50, W~12.30 (no text content)
      * a text shape on top at T~6.70, L~0.70, W~12.0 with the strip text
    """
    spTree = slide.shapes._spTree
    removed = 0
    for sh in list(slide.shapes):
        try:
            top = Emu(sh.top).inches
            left = Emu(sh.left).inches
            width = Emu(sh.width).inches
        except Exception:
            continue

        # The blue band rectangle
        if abs(top - 6.65) < 0.06 and abs(left - 0.50) < 0.06 and width > 12.0:
            if not sh.has_text_frame or not sh.text_frame.text.strip():
                spTree.remove(sh._element)
                removed += 1
                continue

        # The text on top of the band
        if abs(top - 6.70) < 0.06 and abs(left - 0.70) < 0.06 and sh.has_text_frame:
            t = sh.text_frame.text
            if "Guardrail elevates" in t or "Trajectory" in t:
                spTree.remove(sh._element)
                removed += 1

    return removed


# -------------------------------------------------------------- pipeline
def add_pipeline(slide):
    """Render the 5-stage end-to-end pipeline at the bottom of the slide."""
    # Section header
    add_text(slide, 0.85, 6.27, 12.0, 0.22,
             "End-to-end pipeline running in production",
             size=11, bold=True, color=BLUE)

    stages = [
        ("Data Preprocess",        "vitals + nurse notes"),
        ("Scores",                 "qSOFA  \u2022  SIRS  \u2022  SOFA"),
        ("LLM",                    "Claude Sonnet 4.5  \u2022  T = 0"),
        ("LLM-aware Guardrails",   "respects LLM reasoning"),
        ("Noisy-alert Suppression", "blocks non-sepsis patterns"),
    ]

    pipe_top = 6.52
    pipe_h = 0.55
    n = len(stages)
    arrow_w = 0.18
    total_arrow = (n - 1) * arrow_w
    avail = 12.30 - total_arrow
    box_w = avail / n  # ~2.32"

    x = 0.50
    for i, (title, sub) in enumerate(stages):
        add_rounded_box(slide, x, pipe_top, box_w, pipe_h,
                        STAGE_FILL[i], line=STAGE_BORDER[i])
        # Title (bold, accent colour, centered)
        add_text(slide, x + 0.04, pipe_top + 0.06, box_w - 0.08, 0.22,
                 title, size=10, bold=True, color=STAGE_BORDER[i],
                 align=PP_ALIGN.CENTER)
        # Sub (small slate, centered)
        add_text(slide, x + 0.04, pipe_top + 0.30, box_w - 0.08, 0.20,
                 sub, size=8, color=DARK, align=PP_ALIGN.CENTER)
        x += box_w

        # Arrow between stages (except after the last)
        if i < n - 1:
            add_text(slide, x, pipe_top, arrow_w, pipe_h,
                     "\u2192", size=16, bold=True, color=SLATE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            x += arrow_w


# ------------------------------------------------------------------ main
def main():
    if not DECK.exists():
        raise SystemExit(f"Deck not found: {DECK}")

    prs = Presentation(str(DECK))
    if not prs.slides:
        raise SystemExit("Deck has no slides")
    slide = prs.slides[-1]

    n_owned = remove_owned_shapes(slide)
    if n_owned:
        print(f"Removed {n_owned} previously-added shapes (idempotent).")

    n_strip = remove_old_strip(slide)
    if n_strip:
        print(f"Removed old bottom strip ({n_strip} shape(s)).")

    add_pipeline(slide)
    print("Added 5-stage end-to-end pipeline at the bottom.")

    prs.save(str(DECK))
    print(f"\nSaved: {DECK}")


if __name__ == "__main__":
    main()
