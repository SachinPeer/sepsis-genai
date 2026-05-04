import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# === COLOR PALETTE ===
SLATE_900 = RGBColor(15, 23, 42)
SLATE_800 = RGBColor(30, 41, 59)
SLATE_600 = RGBColor(71, 85, 105)
SLATE_500 = RGBColor(100, 116, 139)
SLATE_400 = RGBColor(148, 163, 184)
SLATE_200 = RGBColor(226, 232, 240)
SLATE_100 = RGBColor(241, 245, 249)
WHITE = RGBColor(255, 255, 255)
PURPLE_600 = RGBColor(147, 51, 234)
PURPLE_100 = RGBColor(243, 232, 255)
BLUE_600 = RGBColor(37, 99, 235)
BLUE_500 = RGBColor(59, 130, 246)
BLUE_400 = RGBColor(96, 165, 250)
BLUE_100 = RGBColor(219, 234, 254)
BLUE_50 = RGBColor(239, 246, 255)
GREEN_600 = RGBColor(22, 163, 74)
GREEN_100 = RGBColor(220, 252, 231)
GREEN_50 = RGBColor(240, 253, 244)
GREEN_700 = RGBColor(21, 128, 61)
GREEN_800 = RGBColor(22, 101, 52)
RED_500 = RGBColor(239, 68, 68)
ORANGE_300 = RGBColor(253, 186, 116)
ORANGE_50 = RGBColor(255, 247, 237)
ORANGE_700 = RGBColor(194, 65, 12)


def add_text(slide, left, top, width, height, text, font_size=11,
             bold=False, italic=False, color=SLATE_600, align=PP_ALIGN.LEFT,
             font_name='Arial'):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tb


def add_bullet_list(slide, left, top, width, height, items, font_size=10,
                    color=SLATE_600, bullet="•"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Arial'
        p.space_before = Pt(3)
    return tb


def add_box(slide, left, top, width, height, fill_color=WHITE,
            border_color=SLATE_200, border_width=1):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(border_width)
    return shape


def create_medbeacon_slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(248, 250, 252)

    # ================================================================
    # HEADER
    # ================================================================
    add_text(slide, Inches(0.5), Inches(0.2), Inches(8), Inches(0.5),
             "Medbeacon Architecture", font_size=26, bold=True, color=SLATE_900)
    add_text(slide, Inches(0.5), Inches(0.65), Inches(8), Inches(0.3),
             "Real-Time Clinical Intelligence Workflow & Security Governance",
             font_size=13, italic=True, color=SLATE_500)

    # Confidential tag
    conf = add_box(slide, Inches(11.5), Inches(0.25), Inches(1.5), Inches(0.35),
                   fill_color=SLATE_900, border_color=SLATE_900)
    add_text(slide, Inches(11.5), Inches(0.27), Inches(1.5), Inches(0.3),
             "CONFIDENTIAL", font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0),
                                  Inches(12.3), Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = SLATE_200
    line.line.fill.background()

    # ================================================================
    # ZONE 1: HOSPITAL SOURCE (Left)
    # ================================================================
    Z1_LEFT = Inches(0.4)
    Z1_TOP = Inches(1.3)
    Z1_W = Inches(3.0)
    Z1_H = Inches(3.2)

    # Zone box
    add_box(slide, Z1_LEFT, Z1_TOP, Z1_W, Z1_H, fill_color=WHITE,
            border_color=RGBColor(216, 180, 254), border_width=2)

    # Zone header bar
    z1_hdr = add_box(slide, Z1_LEFT, Z1_TOP, Z1_W, Inches(0.35),
                     fill_color=PURPLE_600, border_color=PURPLE_600)
    add_text(slide, Z1_LEFT + Inches(0.15), Z1_TOP + Inches(0.03), Z1_W, Inches(0.3),
             "ZONE 1: HOSPITAL SOURCE", font_size=9, bold=True, color=WHITE)

    # Icon circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Z1_LEFT + Inches(0.2),
                                    Z1_TOP + Inches(0.55), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PURPLE_100
    circle.line.fill.background()
    add_text(slide, Z1_LEFT + Inches(0.2), Z1_TOP + Inches(0.58), Inches(0.5), Inches(0.45),
             "H", font_size=18, bold=True, color=PURPLE_600, align=PP_ALIGN.CENTER)

    add_text(slide, Z1_LEFT + Inches(0.85), Z1_TOP + Inches(0.55), Inches(2), Inches(0.25),
             "EHR Environment", font_size=12, bold=True, color=SLATE_800)
    add_text(slide, Z1_LEFT + Inches(0.85), Z1_TOP + Inches(0.78), Inches(2), Inches(0.2),
             "Epic / Cerner / Meditech", font_size=9, color=SLATE_500)

    add_bullet_list(slide, Z1_LEFT + Inches(0.25), Z1_TOP + Inches(1.15),
                    Z1_W - Inches(0.5), Inches(1.2),
                    ["Vitals (HR, BP, SpO2, Temp)",
                     "Lab Results (Lactate, WBC, Creatinine)",
                     "Clinician & Nurse Notes",
                     "Patient History & Medications"])

    # Separator line
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Z1_LEFT + Inches(0.2),
                                 Z1_TOP + Inches(2.5), Z1_W - Inches(0.4), Pt(1))
    sep.fill.solid()
    sep.fill.fore_color.rgb = SLATE_100
    sep.line.fill.background()

    add_text(slide, Z1_LEFT + Inches(0.2), Z1_TOP + Inches(2.6), Z1_W - Inches(0.4), Inches(0.3),
             "Data egress via Red Rover EHR Integration Hub",
             font_size=8, italic=True, color=SLATE_400)

    # ================================================================
    # CONNECTOR ARROW
    # ================================================================
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   Inches(3.5), Inches(2.75), Inches(0.6), Inches(0.3))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = BLUE_500
    arrow.line.fill.background()

    # ================================================================
    # ZONE 2: MEDBEACON HIPAA BOUNDARY (Center-Right)
    # ================================================================
    Z2_LEFT = Inches(4.2)
    Z2_TOP = Inches(1.2)
    Z2_W = Inches(8.7)
    Z2_H = Inches(3.4)

    # Dashed boundary
    boundary = add_box(slide, Z2_LEFT, Z2_TOP, Z2_W, Z2_H,
                       fill_color=RGBColor(248, 251, 255),
                       border_color=BLUE_400, border_width=3)
    boundary.line.dash_style = 2

    # Floating label
    label_w = Inches(4.8)
    label_x = Z2_LEFT + (Z2_W - label_w) / 2
    lbl = add_box(slide, label_x, Z2_TOP - Inches(0.18), label_w, Inches(0.35),
                  fill_color=BLUE_600, border_color=BLUE_600)
    add_text(slide, label_x, Z2_TOP - Inches(0.16), label_w, Inches(0.3),
             "MEDBEACON HIPAA SECURE BOUNDARY (AWS BAA + Atlas BAA Covered)",
             font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # --- Card 1: AI Intelligence Engine ---
    C1_LEFT = Z2_LEFT + Inches(0.3)
    C1_TOP = Z2_TOP + Inches(0.45)
    C1_W = Inches(2.9)
    C1_H = Inches(2.7)

    add_box(slide, C1_LEFT, C1_TOP, C1_W, C1_H, fill_color=WHITE,
            border_color=BLUE_100)

    add_text(slide, C1_LEFT + Inches(0.2), C1_TOP + Inches(0.1), C1_W, Inches(0.25),
             "AI INTELLIGENCE ENGINE", font_size=9, bold=True, color=BLUE_600)

    # Sub-card: Predictive Modeling
    sc1 = add_box(slide, C1_LEFT + Inches(0.1), C1_TOP + Inches(0.45),
                  C1_W - Inches(0.2), Inches(0.85),
                  fill_color=RGBColor(248, 250, 252), border_color=SLATE_100)
    add_text(slide, C1_LEFT + Inches(0.2), C1_TOP + Inches(0.5),
             C1_W - Inches(0.4), Inches(0.22),
             "Predictive Modeling", font_size=10, bold=True, color=SLATE_800)
    add_text(slide, C1_LEFT + Inches(0.2), C1_TOP + Inches(0.72),
             C1_W - Inches(0.4), Inches(0.55),
             "Clinical data enriched through SME-defined trend analysis and narrative "
             "preprocessing, then passed to GenAI (AWS Bedrock) for deep reasoning.",
             font_size=8, color=SLATE_500)

    # Sub-card: Narrative Rationale
    sc2 = add_box(slide, C1_LEFT + Inches(0.1), C1_TOP + Inches(1.45),
                  C1_W - Inches(0.2), Inches(0.85),
                  fill_color=RGBColor(248, 250, 252), border_color=SLATE_100)
    add_text(slide, C1_LEFT + Inches(0.2), C1_TOP + Inches(1.5),
             C1_W - Inches(0.4), Inches(0.22),
             "Narrative Rationale", font_size=10, bold=True, color=SLATE_800)
    add_text(slide, C1_LEFT + Inches(0.2), C1_TOP + Inches(1.72),
             C1_W - Inches(0.4), Inches(0.55),
             "Not just a score: provides \"Clinical Rationale\" explaining the "
             "'why' behind every alert to increase clinician trust.",
             font_size=8, color=SLATE_500)

    # --- Card 2: MongoDB Atlas Data Store ---
    DB_LEFT = Z2_LEFT + Inches(3.4)
    DB_TOP = Z2_TOP + Inches(0.45)
    DB_W = Inches(2.5)
    DB_H = Inches(2.7)

    AMBER_50 = RGBColor(255, 251, 235)
    AMBER_200 = RGBColor(253, 230, 138)
    AMBER_700 = RGBColor(180, 83, 9)
    AMBER_800 = RGBColor(146, 64, 14)

    add_box(slide, DB_LEFT, DB_TOP, DB_W, DB_H, fill_color=WHITE,
            border_color=AMBER_200)

    add_text(slide, DB_LEFT + Inches(0.15), DB_TOP + Inches(0.1), DB_W, Inches(0.25),
             "MONGODB ATLAS (on AWS)", font_size=9, bold=True, color=AMBER_700)

    # What we store
    db_store = add_box(slide, DB_LEFT + Inches(0.1), DB_TOP + Inches(0.45),
                       DB_W - Inches(0.2), Inches(1.15),
                       fill_color=AMBER_50, border_color=AMBER_200)
    add_text(slide, DB_LEFT + Inches(0.18), DB_TOP + Inches(0.48),
             DB_W - Inches(0.36), Inches(0.18),
             "Stored (Encrypted AES-256):", font_size=8, bold=True, color=AMBER_800)
    add_bullet_list(slide, DB_LEFT + Inches(0.18), DB_TOP + Inches(0.66),
                    DB_W - Inches(0.36), Inches(0.9),
                    ["Patient demographics (temp, till discharge)",
                     "Vitals, labs, notes (temp, till discharge)",
                     "Doctor feedback (long-term)",
                     "Guardrail configs & audit logs"],
                    font_size=7.5, color=SLATE_600)

    # Safeguards
    db_no = add_box(slide, DB_LEFT + Inches(0.1), DB_TOP + Inches(1.75),
                    DB_W - Inches(0.2), Inches(0.75),
                    fill_color=GREEN_50, border_color=GREEN_100)
    add_text(slide, DB_LEFT + Inches(0.18), DB_TOP + Inches(1.78),
             DB_W - Inches(0.36), Inches(0.18),
             "PHI Safeguards:", font_size=8, bold=True, color=GREEN_800)
    add_bullet_list(slide, DB_LEFT + Inches(0.18), DB_TOP + Inches(1.96),
                    DB_W - Inches(0.36), Inches(0.5),
                    ["Auto-purged on discharge",
                     "LLM never stores patient data"],
                    font_size=7.5, color=SLATE_600)

    # --- Card 3: Clinician Insights ---
    C2_LEFT = Z2_LEFT + Inches(6.1)
    C2_TOP = Z2_TOP + Inches(0.45)
    C2_W = Inches(2.3)
    C2_H = Inches(2.7)

    add_box(slide, C2_LEFT, C2_TOP, C2_W, C2_H, fill_color=WHITE,
            border_color=BLUE_100)

    add_text(slide, C2_LEFT + Inches(0.2), C2_TOP + Inches(0.1), C2_W, Inches(0.25),
             "CLINICIAN INSIGHTS", font_size=9, bold=True, color=GREEN_600)

    # Insight items
    insights = [
        ("Early Warning System", RED_500),
        ("Risk Score + Confidence", BLUE_500),
        ("Guardrail Overrides", GREEN_600),
        ("qSOFA / SIRS / SOFA", BLUE_400),
        ("History-Aware Context", PURPLE_600),
    ]
    for idx, (label, dot_color) in enumerate(insights):
        y = C2_TOP + Inches(0.45) + Inches(idx * 0.26)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     C2_LEFT + Inches(0.15), y + Inches(0.03),
                                     Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = dot_color
        dot.line.fill.background()
        add_text(slide, C2_LEFT + Inches(0.3), y, C2_W - Inches(0.45), Inches(0.22),
                 label, font_size=9, color=SLATE_800)

    # Impact result box
    imp = add_box(slide, C2_LEFT + Inches(0.1), C2_TOP + Inches(1.85),
                  C2_W - Inches(0.2), Inches(0.7),
                  fill_color=GREEN_50, border_color=GREEN_100)
    add_text(slide, C2_LEFT + Inches(0.15), C2_TOP + Inches(1.9),
             C2_W - Inches(0.3), Inches(0.18),
             "IMPACT:", font_size=8, bold=True, color=GREEN_800)
    add_text(slide, C2_LEFT + Inches(0.15), C2_TOP + Inches(2.1),
             C2_W - Inches(0.3), Inches(0.4),
             "360\u00b0 clinical view. Works even when LLM is unavailable.",
             font_size=8, italic=True, color=GREEN_700)

    # ================================================================
    # BOTTOM ROW: COMPLIANCE + OUTCOMES
    # ================================================================
    BOT_TOP = Inches(4.85)

    # --- Compliance Section (Left) ---
    add_box(slide, Inches(0.4), BOT_TOP, Inches(6.3), Inches(2.35),
            fill_color=WHITE, border_color=SLATE_200)

    add_text(slide, Inches(0.65), BOT_TOP + Inches(0.1), Inches(5), Inches(0.3),
             "Why We Are Compliance-First", font_size=14, bold=True, color=SLATE_900)

    compliance_items = [
        ("Temporary Data Only:", "Patient clinical data retained only till discharge, then auto-purged. Encrypted AES-256."),
        ("Zero-Training:", "Amazon Bedrock does NOT store or train on hospital patient data."),
        ("Full Data Isolation:", "PHI never leaves the AWS boundary. MongoDB Atlas hosted on AWS us-east-1."),
        ("Dual BAA Coverage:", "AWS BAA covers Bedrock, EKS, CloudWatch. MongoDB Atlas BAA covers database."),
        ("No PHI in Logs:", "Audit logs contain request IDs and scores only \u2014 zero patient data."),
        ("Certified Trust:", "SOC 2 & ISO 27001 infra. Role-based access control on all data stores."),
    ]

    for idx, (title, desc) in enumerate(compliance_items):
        col = idx % 2
        row = idx // 2
        x = Inches(0.65) + col * Inches(3.1)
        y = BOT_TOP + Inches(0.5) + row * Inches(0.55)

        check = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.03),
                                       Inches(0.16), Inches(0.16))
        check.fill.solid()
        check.fill.fore_color.rgb = GREEN_100
        check.line.fill.background()
        add_text(slide, x + Inches(0.01), y + Inches(0.01), Inches(0.16), Inches(0.16),
                 "\u2713", font_size=8, bold=True, color=GREEN_600, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.22), y, Inches(2.8), Inches(0.45),
                 f"{title} {desc}", font_size=8, color=SLATE_600)

    # --- Outcomes Section (Right, dark) ---
    dark_box = add_box(slide, Inches(6.9), BOT_TOP, Inches(6.1), Inches(2.35),
                       fill_color=SLATE_900, border_color=SLATE_900)

    add_text(slide, Inches(7.15), BOT_TOP + Inches(0.1), Inches(5), Inches(0.3),
             "Powerful Clinical Outcomes", font_size=14, bold=True, color=WHITE)

    add_text(slide, Inches(7.15), BOT_TOP + Inches(0.42), Inches(5.5), Inches(0.25),
             "Medbeacon transitions healthcare from Reactive to Proactive.",
             font_size=10, italic=True, color=SLATE_400)

    outcomes = [
        ("01", "Sepsis Detection",
         "Targeting the #1 cause of hospital mortality with hours-ahead warning."),
        ("02", "Scalable Expansion",
         "Live for Sepsis; Expandable to AKI, DVT/PE, Cardiac Arrest, and Respiratory Failure."),
        ("03", "Clinician Efficiency",
         "Reduces \"Alert Fatigue\" by providing deterministic scores with AI reasoning."),
    ]

    for idx, (num, title, desc) in enumerate(outcomes):
        y = BOT_TOP + Inches(0.82) + idx * Inches(0.48)
        num_box = add_box(slide, Inches(7.15), y, Inches(0.4), Inches(0.4),
                          fill_color=RGBColor(30, 41, 59), border_color=RGBColor(51, 65, 85))
        add_text(slide, Inches(7.15), y + Inches(0.05), Inches(0.4), Inches(0.3),
                 num, font_size=12, bold=True, color=BLUE_400, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(7.7), y, Inches(5), Inches(0.2),
                 title, font_size=9, bold=True, color=SLATE_400)
        add_text(slide, Inches(7.7), y + Inches(0.2), Inches(5), Inches(0.25),
                 desc, font_size=8, italic=True, color=RGBColor(203, 213, 225))

    # ================================================================
    # FOOTER
    # ================================================================
    add_text(slide, Inches(0.5), Inches(7.25), Inches(12.3), Inches(0.25),
             "Medbeacon Inc.  |  www.medbeacon.in  |  AI Clinical Intelligence Platform",
             font_size=8, color=SLATE_400, align=PP_ALIGN.CENTER)

    # Save
    file_name = "docs/Medbeacon_Executive_Overview.pptx"
    prs.save(file_name)
    print(f"PowerPoint created: {file_name}")


if __name__ == "__main__":
    create_medbeacon_slide()
