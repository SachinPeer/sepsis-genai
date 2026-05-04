import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_labeled_icon(slide, left, top, text, color):
    """Helper to add a small circular icon placeholder with text"""
    size = Inches(0.4)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(255, 255, 255)

    tx = slide.shapes.add_textbox(left, top + size, size, Inches(0.3))
    p = tx.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.alignment = PP_ALIGN.CENTER

def create_rich_storyboard():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Colors
    ORANGE = RGBColor(249, 115, 22)
    BLUE = RGBColor(59, 130, 246)
    PURPLE = RGBColor(168, 85, 247)
    RED = RGBColor(239, 68, 68)
    GREEN = RGBColor(34, 197, 94)
    SLATE = RGBColor(15, 23, 42)

    # --- Header ---
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(4), Inches(0.5))
    p = header_box.text_frame.paragraphs[0]
    p.text = "Medbeacon | The Clinical Journey"
    p.font.bold = True
    p.font.size = Pt(18)

    conf_box = slide.shapes.add_textbox(Inches(11.5), Inches(0.1), Inches(1.5), Inches(0.5))
    p = conf_box.text_frame.paragraphs[0]
    p.text = "CONFIDENTIAL"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = ORANGE

    # --- Timeline Step Data ---
    steps = [
        {"time": "5:55 AM", "title": "SYSTEM WAKES UP", "desc": "Scans vitals, labs, and notes every 5 mins.", "color": ORANGE, "bg": RGBColor(255, 251, 235)},
        {"time": "6:01 AM", "title": "UNCOVERING PATTERN", "desc": "Heart change (105) & BP change (98).", "color": BLUE, "bg": RGBColor(239, 246, 255)},
        {"time": "6:01 AM", "title": "INTELLIGENT ANALYSIS", "desc": "GenAI finds truth in nurse concerns.", "color": PURPLE, "bg": RGBColor(245, 243, 255)},
        {"time": "6:02 AM", "title": "HIGH-PRIORITY ALERT", "desc": "RISK: 72 (HIGH). Bile leak risk flagged.", "color": RED, "bg": RGBColor(254, 242, 242)},
        {"time": "6:05 AM", "title": "DR. PRIYA REVIEWS", "desc": "Human-in-the-loop review of AI rationale.", "color": BLUE, "bg": RGBColor(248, 250, 252)},
        {"time": "6:15 AM", "title": "EARLY INTERVENTION", "desc": "Orders antibiotics. Hours before crash.", "color": RED, "bg": RGBColor(255, 241, 242)},
        {"time": "12:00 PM", "title": "PATIENT STABLE", "desc": "Early intervention prevented shock.", "color": GREEN, "bg": RGBColor(240, 253, 244)}
    ]

    left_m = Inches(0.3)
    top_m = Inches(0.8)
    card_w = Inches(1.75)
    card_h = Inches(3.8)
    spacing = Inches(1.82)

    for i, step in enumerate(steps):
        x = left_m + (i * spacing)

        # Step Card
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top_m, card_w, card_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = step["bg"]
        rect.line.color.rgb = step["color"]
        rect.line.width = Pt(1.5)

        # Time Header
        time_rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.3), top_m - Inches(0.15), Inches(1.1), Inches(0.35))
        time_rect.fill.solid()
        time_rect.fill.fore_color.rgb = RGBColor(255, 255, 255)
        time_rect.line.color.rgb = step["color"]

        t_tf = time_rect.text_frame
        p = t_tf.paragraphs[0]
        p.text = step["time"]
        p.font.size = Pt(10)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Body Text
        txt = slide.shapes.add_textbox(x + Inches(0.1), top_m + Inches(1.8), card_w - Inches(0.2), Inches(1.5))
        txt.text_frame.word_wrap = True
        p1 = txt.text_frame.paragraphs[0]
        p1.text = step["title"]
        p1.font.bold = True
        p1.font.size = Pt(11)
        p1.alignment = PP_ALIGN.CENTER

        p2 = txt.text_frame.add_paragraph()
        p2.text = step["desc"]
        p2.font.size = Pt(9)
        p2.alignment = PP_ALIGN.CENTER

    # --- DOCTOR FEEDBACK LOOP ---
    loop_x, loop_y = Inches(7.8), Inches(2.2)
    loop_shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, loop_x, loop_y, Inches(1.5), Inches(0.5))
    loop_shape.rotation = 90
    loop_shape.fill.solid()
    loop_shape.fill.fore_color.rgb = PURPLE

    loop_txt = slide.shapes.add_textbox(loop_x + Inches(0.6), loop_y - Inches(0.5), Inches(1.5), Inches(0.5))
    p = loop_txt.text_frame.paragraphs[0]
    p.text = "DOCTOR FEEDBACK"
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    # --- BEHIND THE SCENES WORKFLOW (Bottom) ---
    bt_y = Inches(5.0)
    bt_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), bt_y, Inches(12.7), Inches(2.2))
    bt_bg.fill.solid()
    bt_bg.fill.fore_color.rgb = SLATE

    bt_title = slide.shapes.add_textbox(Inches(0.5), bt_y + Inches(0.1), Inches(5), Inches(0.4))
    p = bt_title.text_frame.paragraphs[0]
    p.text = "BEHIND THE SCENES END-TO-END DATA FLOW"
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.size = Pt(12)
    p.font.bold = True

    # Workflow Blocks
    nodes = ["1. Monitors", "2. Hospital EHR", "3. Red Rover API", "4. Backend EKS", "5. Preprocessor", "6. GenAI", "7. Guardrails", "8. Dashboard"]
    n_w, n_h = Inches(1.4), Inches(0.6)
    n_x_start = Inches(0.5)
    n_spacing = Inches(1.55)

    for i, node in enumerate(nodes):
        nx = n_x_start + (i * n_spacing)
        n_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, nx, bt_y + Inches(0.8), n_w, n_h)
        n_rect.fill.solid()
        n_rect.fill.fore_color.rgb = RGBColor(51, 65, 85)
        n_rect.line.color.rgb = BLUE

        p = n_rect.text_frame.paragraphs[0]
        p.text = node
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # Latency Bar
    lat_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), bt_y + Inches(1.7), Inches(5), Inches(0.3))
    lat_bar.fill.solid()
    lat_bar.fill.fore_color.rgb = RGBColor(71, 85, 105)
    p = lat_bar.text_frame.paragraphs[0]
    p.text = "TOTAL: ~10 SECONDS FROM DATA TO ALERT"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    prs.save("presentations/current/Medbeacon_Clinical_Storyboard.pptx")
    print("PowerPoint created: presentations/current/Medbeacon_Clinical_Storyboard.pptx")

if __name__ == "__main__":
    create_rich_storyboard()
