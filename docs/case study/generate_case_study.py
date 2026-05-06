"""
Generate Medbeacon Sepsis GenAI one-pager case study.

Based on the Winter masterdeck template, specifically slide 6 ("The Impact"),
which has a 3-row left column (Challenge / Solution / Impact) + 2x2 KPI grid
on the right. All Winter fonts (Segoe UI family) and color palette
(#454545 / #898A87 / #506370) are preserved by editing in place.

Output: docs/case study/Medbeacon_Sepsis_CaseStudy.pptx
"""
import shutil
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = Path(__file__).parent
TEMPLATE = HERE / "Winter masterdeck template.pptx"
OUTPUT = HERE / "Medbeacon_Sepsis_CaseStudy.pptx"

KEEP_SLIDE_INDEX = 5  # 0-based; slide 6 "The Impact"

# --- Winter palette (confirmed from slide 6 inspection) ---
COLOR_TITLE = RGBColor(0x45, 0x45, 0x45)   # #454545
COLOR_BODY = RGBColor(0x89, 0x8A, 0x87)    # #898A87
COLOR_LABEL = RGBColor(0x50, 0x63, 0x70)   # #506370

FONT_REGULAR = "Segoe UI"
FONT_SEMIBOLD = "Segoe UI Semibold"


def set_text(tf, text, font_name, size_pt, *, bold=False, color=None):
    """Replace text in a text_frame, preserving single-paragraph styling."""
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def add_trust_strip(slide):
    """Add a horizontal strip below the 2x2 cards showing the 3-step pipeline
    plus HIPAA + Auditable trust badges. Strip sits at T~6.75, spans the
    right panel width.
    """
    # Strip occupies the right-panel width (matches the existing grey panel
    # which extends from L~4.27 to L~12.65).
    strip_left = Inches(4.27)
    strip_top = Inches(6.78)
    strip_w = Inches(8.38)
    strip_h = Inches(0.42)

    # Background pill with very subtle fill (slightly darker than the panel grey)
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, strip_left, strip_top, strip_w, strip_h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF0)  # very light grey
    bg.line.fill.background()  # no border
    bg.adjustments[0] = 0.5  # rounder corners

    # Single text line: "Preprocess → AI Reasoning → Guardrails    ·    HIPAA-compliant    ·    Fully auditable"
    tb = slide.shapes.add_textbox(strip_left, strip_top, strip_w, strip_h)
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    # Mixed-style runs in one paragraph: stages bold, dot separators muted, badges semibold accent
    def add_run(text, *, bold=False, semibold=False, color=COLOR_BODY, size=11):
        r = p.add_run()
        r.text = text
        r.font.name = FONT_SEMIBOLD if semibold else FONT_REGULAR
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color

    add_run("3-stage pipeline: ", semibold=True, color=COLOR_LABEL, size=10)
    add_run("Preprocess  ", color=COLOR_TITLE, size=11, bold=True)
    add_run("→  ", color=COLOR_BODY, size=11)
    add_run("AI Reasoning  ", color=COLOR_TITLE, size=11, bold=True)
    add_run("→  ", color=COLOR_BODY, size=11)
    add_run("Guardrails", color=COLOR_TITLE, size=11, bold=True)
    add_run("        ", color=COLOR_BODY, size=11)
    add_run("HIPAA-compliant", semibold=True, color=COLOR_LABEL, size=11)
    add_run("    ·    ", color=COLOR_BODY, size=11)
    add_run("Fully auditable", semibold=True, color=COLOR_LABEL, size=11)


def remove_all_slides_except(prs, keep_index):
    """Remove every slide except `keep_index` from the presentation."""
    sldIdLst = prs.slides._sldIdLst
    slide_ids = list(sldIdLst)
    # Map slide id entries to their index
    # Must remove both the sldId reference and the related slide part.
    for i, sldId in enumerate(slide_ids):
        if i == keep_index:
            continue
        rId = sldId.get(qn('r:id'))
        try:
            prs.part.drop_rel(rId)
        except KeyError:
            pass
        sldIdLst.remove(sldId)


def build():
    if not TEMPLATE.exists():
        raise FileNotFoundError(f"Template not found: {TEMPLATE}")
    shutil.copy(TEMPLATE, OUTPUT)

    prs = Presentation(str(OUTPUT))
    remove_all_slides_except(prs, KEEP_SLIDE_INDEX)
    assert len(prs.slides) == 1, f"Expected 1 slide, got {len(prs.slides)}"
    slide = prs.slides[0]

    # Text we need to replace, keyed by the EXACT original text.
    # Same-text collisions (multiple shapes with identical placeholder text)
    # are disambiguated by position via a secondary rule below.
    replacements = {
        # --- Header area ---
        "THE IMPACT": ("CASE STUDY", FONT_SEMIBOLD, 10, False, COLOR_TITLE),
        "Case title": ("Medbeacon Sepsis GenAI", FONT_REGULAR, 22, False, COLOR_TITLE),
        "Cases across industries": ("Capability snapshot · 2026", FONT_REGULAR, 8, False, COLOR_BODY),

        # --- Left column headers ---
        "The challenge": ("The challenge", FONT_REGULAR, 15, False, COLOR_TITLE),
        "The solution": ("The solution", FONT_REGULAR, 15, False, COLOR_TITLE),
        "The impact": ("The impact", FONT_REGULAR, 15, False, COLOR_TITLE),

        # --- 2x2 KPI cards: category labels (top) ---
        "Healthcare": ("Nurse Notes", FONT_SEMIBOLD, 14, False, COLOR_LABEL),
        "Financial services": ("Clinical Scoring", FONT_SEMIBOLD, 14, False, COLOR_LABEL),
        "Transportation": ("Guardrails", FONT_SEMIBOLD, 14, False, COLOR_LABEL),
        "Automotive": ("Feedback Loop", FONT_SEMIBOLD, 14, False, COLOR_LABEL),

        # --- 2x2 KPI cards: big feature headlines ---
        "$12m": ("Unstructured", FONT_REGULAR, 24, False, COLOR_TITLE),
        "$8m": ("Deterministic", FONT_REGULAR, 24, False, COLOR_TITLE),
        "$3m": ("Per-Hospital", FONT_REGULAR, 24, False, COLOR_TITLE),
        "$120m": ("Closed-Loop", FONT_REGULAR, 24, False, COLOR_TITLE),

        # --- 2x2 KPI cards: caption under each (keep ~35 chars to stay on one line) ---
        "Saved by improving the revenue cycle":
            ("Free-text nurse notes understood", FONT_SEMIBOLD, 14, False, COLOR_LABEL),
        "Saved by accelerating global payments":
            ("qSOFA · SIRS · SOFA, baked in", FONT_SEMIBOLD, 14, False, COLOR_LABEL),
        "Saved by real time flight operations":
            ("Configurable, auditable, hot-reload", FONT_SEMIBOLD, 14, False, COLOR_LABEL),
        "Saved by accelerating manufacturing":
            ("Clinician feedback · continuous tuning", FONT_SEMIBOLD, 14, False, COLOR_LABEL),

        # --- Page number ---
        "01": ("01", FONT_REGULAR, 9, False, COLOR_TITLE),
    }

    # The Lorem Ipsum body shapes and the large description shape all have
    # identical placeholder text, so we disambiguate them by position (T, L).
    # From Winter slide 6 inspection:
    #   L=0.26 T=1.63 — header subtitle (Segoe UI 11pt #898A87)
    #   L=0.34 T=3.23 — challenge body (Segoe UI 11pt #898A87)
    #   L=0.34 T=4.60 — solution body (Segoe UI 11pt #898A87)
    #   L=0.34 T=5.99 — impact body (Segoe UI 11pt #898A87)
    #   L=4.46 T=2.93 — card 1 description (Segoe UI 11pt)
    #   L=8.78 T=2.93 — card 2 description (Segoe UI 11pt)
    #   L=4.46 T=5.85 — card 3 description (Segoe UI 11pt)
    #   L=8.79 T=5.82 — card 4 description (Segoe UI 11pt)
    POS_REPLACEMENTS = {
        # (L_in, T_in) -> new text
        (0.26, 1.63): (
            "A clinical decision-support system for sepsis risk: nurse-note understanding, "
            "deterministic scoring, hospital-configurable guardrails, and a clinician feedback loop.",
            FONT_REGULAR, 11, False, COLOR_BODY,
        ),
        (0.34, 3.23): (
            "Critical sepsis signals — confusion, mottling, lethargy — live in free-text "
            "nurse notes that rule-based EHR alarms ignore. Alert fatigue and AI opacity "
            "both erode clinician trust.",
            FONT_REGULAR, 11, False, COLOR_BODY,
        ),
        (0.34, 4.60): (
            "Medbeacon layers four capabilities: LLM-powered understanding of unstructured "
            "nurse notes, deterministic qSOFA / SIRS / SOFA math, hospital-configurable "
            "safety guardrails, and a clinician-driven feedback loop.",
            FONT_REGULAR, 11, False, COLOR_BODY,
        ),
        (0.34, 5.99): (
            "A clinically safe, auditable, tunable decision-support system — not a black "
            "box. Each hospital owns its thresholds, override rules and calibration. "
            "Validated on 90 ICU patients to prove the architecture works on real data.",
            FONT_REGULAR, 11, False, COLOR_BODY,
        ),
        (4.46, 2.93): (
            "Free-text observations — confusion, mottling, lethargy — parsed by the LLM into "
            "structured sepsis signals that rule-based EHR systems simply cannot see.",
            FONT_REGULAR, 11, False, COLOR_BODY,
        ),
        (8.78, 2.93): (
            "qSOFA, SIRS and SOFA computed deterministically on every evaluation — "
            "audit-ready clinical math that clinicians can trace alongside the AI rationale.",
            FONT_REGULAR, 11, False, COLOR_BODY,
        ),
        (4.46, 5.85): (
            "Each hospital tunes its own thresholds, override rules, and alert levels. "
            "Changes hot-reload without downtime; every decision is logged and auditable.",
            FONT_REGULAR, 11, False, COLOR_BODY,
        ),
        (8.79, 5.82): (
            "Clinician accept / reject feedback drives prompt tuning and threshold "
            "calibration — the system sharpens on your own patients, not generic data.",
            FONT_REGULAR, 11, False, COLOR_BODY,
        ),
    }

    replaced_keys = set()
    replaced_positions = set()

    # Shapes whose width we need to bump up so new text fits on one line.
    # Keyed by ORIGINAL text.
    WIDTH_OVERRIDES = {
        "THE IMPACT": 2.0,     # "CASE STUDY" (wider than "THE IMPACT")
        "Case title": 6.5,     # "Medbeacon Sepsis GenAI"
    }

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        current = shape.text_frame.text.strip()
        L = shape.left / 914400 if shape.left else 0
        T = shape.top / 914400 if shape.top else 0

        # Apply width override BEFORE replacing text
        if current in WIDTH_OVERRIDES:
            shape.width = Inches(WIDTH_OVERRIDES[current])

        # Position-based replacement first (more specific)
        matched_pos = None
        for (pl, pt), payload in POS_REPLACEMENTS.items():
            if abs(L - pl) < 0.1 and abs(T - pt) < 0.1:
                matched_pos = (pl, pt)
                break
        if matched_pos is not None:
            new_text, fn, sz, bold, color = POS_REPLACEMENTS[matched_pos]
            set_text(shape.text_frame, new_text, fn, sz, bold=bold, color=color)
            replaced_positions.add(matched_pos)
            continue

        # Text-based replacement
        if current in replacements:
            new_text, fn, sz, bold, color = replacements[current]
            set_text(shape.text_frame, new_text, fn, sz, bold=bold, color=color)
            replaced_keys.add(current)

    # Report
    missing_keys = set(replacements.keys()) - replaced_keys
    missing_pos = set(POS_REPLACEMENTS.keys()) - replaced_positions
    if missing_keys:
        print(f"WARNING — text keys not found: {sorted(missing_keys)}")
    if missing_pos:
        print(f"WARNING — positions not found: {sorted(missing_pos)}")

    # --- Trust strip below the 2x2 cards: 3-step process + HIPAA + Auditable ---
    add_trust_strip(slide)

    prs.save(str(OUTPUT))
    print(f"Generated: {OUTPUT}")
    print(f"Text replacements: {len(replaced_keys)}/{len(replacements)}")
    print(f"Position replacements: {len(replaced_positions)}/{len(POS_REPLACEMENTS)}")


if __name__ == "__main__":
    build()
